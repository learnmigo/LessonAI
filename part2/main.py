import asyncio
import os
import json
import re
import time
import aiofiles
import boto3
import base64

from uuid import uuid4
from fastapi import FastAPI, UploadFile, File, HTTPException
from openai import OpenAI
from pydantic import BaseModel
from typing import List
from io import BytesIO
from fastapi.responses import FileResponse
from mangum import Mangum
from fastapi.middleware.cors import CORSMiddleware
from urllib.parse import quote_plus
from pptx import Presentation as Presentation_pptx
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
from loguru import logger

app = FastAPI()
handler = Mangum(app)

origins = ["*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

query = """
# Make a PPT presentation and a script which could be said alongside the presentation based on the {sublearning_outcome} The output should be in a JSON FORMAT ONLY WITHOUT ANY EXPLANATORY TEXT AROUND IT which returns the slide title, slide content in bullet points, 2-3 keywords based on the slide content and the script which could be said by the teacher when the current slide is shown. 

# The ppt should have the following slides: 
 
# Introduction slide which will use the hook/gain attention, establish relevance  and  Recall/Activate memory reparts of {sublearning_outcome}

# Objective slide which will contain mindmap {sublearning_outcome}

# Demonstration slide which will contain th demonstration part of topic 

# Practice assessments slide which will contain Practice Assessments 

# Summary slide which contains Summary + Priming for Next Topic:

# - 

# THE SCRIPT SHOULD BE BASED ON THE CONTENT OF EACH SLIDE AND THE EXPLANATION GIVEN IN THE {sublearning_outcome}

# THE RESPONSE SHOULD BE IN THE FOLLOWING JSON FORMAT ONLY. DO NOT ADD ANY EXPLANATORY TEXT BEFORE OR AFTER IT.

# {
#   "slides": [
#     {
#       "title": "string",
#       "content": [
#         "string"
#       ],
#       "keywords": [
#         "string"
#       ],
#       "script": "string"
#     }
#   ]
# }
# EACH SLIDE WOULD BE A LIST ITEM IN THE ABOVE JSON: 
# TITLE - TITLE FOR THE SLIDE CONTENT
# CONTENT - 3 TO 5 BULLET POINTS EXPLAINING THE CONTENT
# KEYWORDS - 2 TO 3 SINGLE WORD KEYWORDS ABOUT THE ENTIRE SLIDE. THESE KEYWORDS WOULD BE USED TO SEARCH FOR IMAGES SO KEEP THEM APT ALSO MAKE SURE KEYWORDS FOR EACH INDIVIDUAL SLIDES MUST BE DIFFERENT, THERE SHOULD NOT BE ANY REPEATING KEYWORDS
# SCRIPT - A COMPREHENSIVE SCRIPT TO GO ALONG WITH THE SLIDE WHICH EXPLAINS THE SLIDE IN DEPTH WITH ELLUSIVE EXAMPLES

# """

ACCESS_KEY = os.getenv("S3_ACCESS_KEY")
SECRET_KEY = os.getenv("S3_SECRET_KEY")
api_key = os.getenv("OPENAI_API_KEY")
assistant_id = os.getenv("OPENAI_ASSISTANT_ID")
AWS_BUCKET_THUMBNAIL = 'ppt-thumbnail-bucket-1'
AWS_BUCKET_PPT = 'ppt-bucket-1'
AWS_BUCKET_IMAGE = 'ppt-image-bucket-1'


client = OpenAI(api_key=api_key)

s3_client = boto3.client(
    's3',
    aws_access_key_id=ACCESS_KEY,
    aws_secret_access_key=SECRET_KEY,
)

session = boto3.Session(
    aws_access_key_id=ACCESS_KEY,
    aws_secret_access_key=SECRET_KEY,
)

s3 = session.resource('s3')
bucketThumbnail = s3.Bucket(AWS_BUCKET_THUMBNAIL)
bucketPPT = s3.Bucket(AWS_BUCKET_PPT)
bucketImage = s3.Bucket(AWS_BUCKET_IMAGE)

async def s3_upload_thumbnail(content: bytes, key: str):
    logger.info(f'Uploading thumbnail {key} to s3')
    await asyncio.to_thread(bucketThumbnail.put_object, Key=key, Body=content)

async def s3_upload_ppt(content: bytes, key: str):
    logger.info(f'Uploading ppt {key} to s3')
    await asyncio.to_thread(bucketPPT.put_object, Key=key, Body=content)

async def s3_upload_image(content: bytes, key: str):
    logger.info(f'Uploading image {key} to s3')
    await asyncio.to_thread(bucketImage.put_object, Key=key, Body=content)

async def s3_download_image_binary(key: str):
    logger.info(f'Downloading Image {key} from s3')
    try:
        response = await asyncio.to_thread(s3_client.get_object, Bucket=AWS_BUCKET_IMAGE, Key=key)
        base64_content = response['Body'].read().decode('utf-8')
        decoded_content = base64.b64decode(base64_content)
    except Exception as e:
        logger.error(f"Error getting object {key} from bucket {AWS_BUCKET_IMAGE}. {e}")
        return None
    return decoded_content

TEMP_DIR = "/tmp"

async def create_thread():
    thread = await asyncio.to_thread(client.beta.threads.create)
    thread_id = thread.id
    logger.info("THREAD CREATED!!!")
    return thread_id

async def run_thread(assistant_id, thread_id):
    response = await asyncio.to_thread(client.beta.threads.runs.create, assistant_id=assistant_id, thread_id=thread_id)
    logger.info("THREAD RUN!!!")
    return response

async def write_msg(query, thread_id):
    thread_message = await asyncio.to_thread(client.beta.threads.messages.create, thread_id=thread_id, role="user", content=query)
    logger.info(thread_message)

class Slide(BaseModel):
    title: str
    content: List[str]
    keywords: List[str]
    script: str
    image_s3: str = ""

class Slides(BaseModel):
    slides: List[Slide]
    template: str

class JSONScriptInput(BaseModel):
    subLearningOutcome: str

def format_text(input_text):
    formatted_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', input_text)
    return formatted_text

async def generate_image(query, index):
    query = f'Generate an Image: {query}'
    logger.info(f"Slide: {index}")
    logger.info(query)
    response = await asyncio.to_thread(client.images.generate, model="dall-e-3", prompt=query, n=1, quality="standard", response_format="b64_json")
    image_data = response.data[0].b64_json
    key = f"{uuid4()}.txt"
    await s3_upload_image(key=key, content=image_data)
    return key

async def create_thumbnail(slide, ppt_file_path):
    logger.info("Creating thumbnail")
    img = Image.new("RGB", (1280, 720), color="white")
    draw = ImageDraw.Draw(img)
    font_file_path = "Gidole-Regular.ttf"
    font_title_size = 50
    font_content_size = 30
    font_title = ImageFont.truetype(font_file_path, font_title_size)
    font_content = ImageFont.truetype(font_file_path, font_content_size)
    draw.text((100, 50), slide["title"], fill="black", font=font_title)
    content_y = 150
    for line in slide["content"]:
        draw.text((100, content_y), "• " + line, fill="black", font=font_content)
        content_y += 50
    image_s3 = slide['image_s3']
    if isinstance(image_s3, str) and image_s3:
        try:
            image_response = await s3_download_image_binary(key=image_s3)
            if image_response:
                image = Image.open(BytesIO(image_response))
                image = image.resize((400, 300))
                img.paste(image, (800, 350))
        except Exception as e:
            logger.error(f"Failed to fetch image: {e}")
    uniqueID = uuid4()
    img_file_name = f"{uniqueID}.jpg"
    img_file_path = os.path.join(TEMP_DIR, img_file_name)
    img.save(img_file_path)
    async with aiofiles.open(img_file_path, mode='rb') as f:
        content = await f.read()
    await s3_upload_thumbnail(content, img_file_name)
    return img_file_name

async def generate_ppt_doc(slide_formatted_data, template):
    fetch_image_tasks = []
    for i, slide in enumerate(slide_formatted_data):
        if 'keywords' in slide:
            keywords = slide['keywords']
            combined_keywords = ' ,'.join(keywords)
            content = slide['content']
            combined_content = ' ,'.join(content)
            dallE_Query = "KeyWords: " + combined_keywords + " Content: " + combined_content
            fetch_image_tasks.append(generate_image(dallE_Query, i))
    image_results = await asyncio.gather(*fetch_image_tasks)
    for slide, image_s3 in zip(slide_formatted_data, image_results):
        if 'keywords' in slide:
            slide['image_s3'] = image_s3
    prs = Presentation_pptx(f'./templates/{template}')
    for slide_data in slide_formatted_data:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_data['title']
        left_content = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(5.5), Inches(4))
        left_text_frame = left_content.text_frame
        left_text_frame.word_wrap = True
        for text in slide_data['content']:
            p = left_text_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(18)
            p.space_after = Pt(12)
            p.level = 0
        image_s3 = slide_data['image_s3']
        if isinstance(image_s3, str) and image_s3:
            try:
                image_response = await s3_download_image_binary(key=image_s3)
                if image_response:
                    image_stream = BytesIO(image_response)
                    left = Inches(7.5)
                    top = Inches(2.2)
                    width = height = Inches(4)
                    slide.shapes.add_picture(image_stream, left, top, width, height)
            except Exception as e:
                logger.error(f"Failed to fetch image: {e}")
    uniqueID = uuid4()
    ppt_file_name = f'{uniqueID}.pptx'
    ppt_file_path = f'/tmp/ppt_{ppt_file_name}'
    prs.save(ppt_file_path)
    thumbnail_s3_file = await create_thumbnail(slide_formatted_data[0], ppt_file_path)
    async with aiofiles.open(ppt_file_path, mode='rb') as f:
        content = await f.read()
    await s3_upload_ppt(content, ppt_file_name)
    return {
        "thumbnail": thumbnail_s3_file,
        "ppt": ppt_file_name,
        "slides": slide_formatted_data
    }

@app.post("/json_script")
async def upload_file(input: JSONScriptInput):
    try:
        sub_learning_outcome = input.subLearningOutcome
        thread_id = await create_thread()
        await write_msg(query + "\n Below is the given Sub-Learning Outcome: \n" + sub_learning_outcome, thread_id)
        response = await run_thread(assistant_id, thread_id)
        i = 1
        while response.status.lower() != "completed":
            response = await asyncio.to_thread(client.beta.threads.runs.retrieve, run_id=response.id, thread_id=thread_id)
            await asyncio.sleep(1)
            i += 1
        query_response = await asyncio.to_thread(client.beta.threads.messages.list, thread_id=thread_id)
        for answer in query_response:
            response_text = answer.content[0].text.value
            break
        start_index = response_text.find('{')
        end_index = response_text.rfind('}') + 1
        json_str = response_text[start_index:end_index]
        json_data = json.loads(json_str)
        return json_data
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate Response: {str(e)}")

@app.post("/generate_ppt")
async def generate_ppt(slides: Slides):
    try:
        slidesArr = []
        for slide in slides.slides:
            slideObj = {
                "title": slide.title,
                "content": slide.content,
                "keywords": slide.keywords,
                "script": slide.script,
                "image_s3": slide.image_s3
            }
            slidesArr.append(slideObj)
        ppt_result = await generate_ppt_doc(slidesArr, slides.template)
        return ppt_result
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PPT: {str(e)}")

@app.get("/")
async def home():
    return {
        "msg": "Hello World!"
    }
